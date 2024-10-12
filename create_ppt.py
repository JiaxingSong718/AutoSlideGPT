from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

def create_ppt(text_file, design_number, ppt_name):
    prs = Presentation(f"Designs/Design-{design_number}.pptx")
    slide_count = 0
    header = ""
    paragraphs = []  # 存储段落标题和描述

    # 读取文本文件并生成内容
    with open(text_file, 'r', encoding='utf-8') as f:
        for line in f:
            if line.startswith('Slide:'):
                # 添加新幻灯片前检查是否需要添加之前的内容
                if slide_count > 0:
                    slide = prs.slides.add_slide(prs.slide_layouts[1])
                    title = slide.shapes.title
                    title.text = header

                    # 判断是否是第一张幻灯片，进行左右分开布局
                    if slide_count == 1:
                        # 左右分开布局
                        # 计算文本框的位置和大小
                        slide_width = prs.slide_width
                        slide_height = prs.slide_height

                        # 左右分开布局
                        left_textbox_width = Inches(4)
                        right_textbox_width = Inches(4)
                        left_textbox_height = Inches(4)
                        right_textbox_height = Inches(4)

                        # 计算左侧和右侧文本框的 X 坐标，以居中
                        left_textbox_x = (slide_width - left_textbox_width - right_textbox_width) / 2
                        right_textbox_x = left_textbox_x + left_textbox_width

                        # 添加文本框，Y 坐标向下移动到整个页面高度的 1/4 + 0.5 英寸
                        text_offset = slide_height / 4 + Inches(1)  # 页面高度的四分之一再向下移动 0.5 英寸
                        left_textbox = slide.shapes.add_textbox(left_textbox_x, text_offset, left_textbox_width, left_textbox_height)
                        right_textbox = slide.shapes.add_textbox(right_textbox_x, text_offset, right_textbox_width, right_textbox_height)
                        left_tf = left_textbox.text_frame
                        right_tf = right_textbox.text_frame

                        # 将段落内容左右布局添加
                        mid_point = len(paragraphs) // 2
                        for i, paragraph in enumerate(paragraphs):
                            # 左侧放前半部分
                            if i < mid_point:
                                left_p = left_tf.add_paragraph()
                                run_title = left_p.add_run()
                                run_title.text = paragraph['title']
                                run_title.font.size = Pt(16)  # 设置标题字体大小
                                run_title.font.bold = True  # 加粗标题
                            # 右侧放后半部分
                            else:
                                right_p = right_tf.add_paragraph()
                                run_title = right_p.add_run()
                                run_title.text = paragraph['title']
                                run_title.font.size = Pt(16)  # 设置标题字体大小
                                run_title.font.bold = True  # 加粗标题
                    else:
                        # 普通布局，将段落内容放入幻灯片的默认文本框中
                        body_shape = slide.shapes.placeholders[1]
                        tf = body_shape.text_frame
                        tf.clear()  # 清空之前的内容
                        for paragraph in paragraphs:
                            p = tf.add_paragraph()
                            # 设置段落标题
                            run_title = p.add_run()
                            run_title.text = paragraph['title'] + ": "
                            run_title.font.size = Pt(16)  # 设置标题字体大小
                            run_title.font.bold = True  # 加粗标题

                            # 插入换行
                            p = tf.add_paragraph()  # 新增段落

                            # 设置描述
                            run_description = p.add_run()
                            run_description.text = paragraph['description']
                            run_description.font.size = Pt(12)  # 设置描述字体大小
                            run_description.font.bold = False  # 描述不加粗

                            p.space_after = Inches(0.1)  # 设置段落后间距

                    paragraphs = []  # 重置段落列表

                slide_count += 1
                header = ""  # 重置header
                continue

            elif line.startswith('Header:'):
                header = line.replace('Header:', '').strip()
                continue
            
            elif line.startswith('Title:'):
                # 提取标题并生成封面页
                title = line.replace('Title:', '').strip()
                cover_slide = prs.slides.add_slide(prs.slide_layouts[0])  # 选择封面布局
                title_shape = cover_slide.shapes.title
                title_shape.text = title  # 设置封面标题

                # 在封面页底部添加 AutoSlideGPT
                left_inch = Inches(5.4)
                top_inch = Inches(5) 
                width_inch = Inches(2.62)
                height_inch = Inches(1)
                textbox = cover_slide.shapes.add_textbox(left_inch, top_inch, width_inch, height_inch)
                tf = textbox.text_frame
                p = tf.add_paragraph()
                run = p.add_run()
                run.text = "汇报人：AutoSlideGPT"
                run.font.size = Pt(18)  # 设置字体大小
                run.font.bold = True  # 设置加粗
                run.font.color.rgb = RGBColor(0, 0, 0)  # 设置字体颜色为黑色
                continue

            elif line.startswith('Paragraph'):
                # 去掉前缀并读取标题
                paragraph_title = line.split(':', 1)[-1].strip().rstrip(',')
                # 读取描述行
                description_line = f.readline().strip()  # 读取描述
                # 将段落标题和描述存入列表
                description = description_line.replace('description:', '').strip()
                paragraphs.append({'title': paragraph_title, 'description': description})

    # 添加最后一个幻灯片（确保有内容）
    if slide_count > 0 and paragraphs:
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        title = slide.shapes.title
        title.text = header

        body_shape = slide.shapes.placeholders[1]
        tf = body_shape.text_frame
        tf.clear()  # 清空之前的内容
        for paragraph in paragraphs:
            p = tf.add_paragraph()
            # 设置段落标题
            run_title = p.add_run()
            run_title.text = paragraph['title'] + ": "
            run_title.font.size = Pt(16)  # 设置标题字体大小
            run_title.font.bold = True  # 加粗标题

            # 插入换行
            p = tf.add_paragraph()  # 新增段落

            # 设置描述
            run_description = p.add_run()
            run_description.text = paragraph['description']
            run_description.font.size = Pt(12)  # 设置描述字体大小
            run_description.font.bold = False  # 描述不加粗

            p.space_after = Inches(0.1)  # 设置段落后间距

    # 保存PPT文件
    prs.save(f'GeneratedPresentations/{ppt_name}.pptx')
    file_path = f"GeneratedPresentations/{ppt_name}.pptx"
    return f"{file_path}"